"""
Tests for converters/__init__.py — convert_file, get_supported_extensions,
and image extraction helpers.

Covers:
- convert_file function with various formats
- get_supported_extensions
- Image extraction from DOCX/PPTX (via _extract_docx_images, _extract_pptx_images)
- FileNotFoundError for missing files
"""

import pytest
import tempfile
import zipfile
from pathlib import Path
from doc_knowledge.converters import (
    convert_file,
    get_supported_extensions,
)


class TestGetSupportedExtensions:
    """获取支持的文件扩展名列表"""

    def test_returns_list(self):
        exts = get_supported_extensions()
        assert isinstance(exts, list)

    def test_contains_common_formats(self):
        exts = get_supported_extensions()
        assert ".pdf" in exts
        assert ".docx" in exts
        assert ".pptx" in exts
        assert ".xlsx" in exts
        assert ".xls" in exts
        assert ".html" in exts
        assert ".txt" in exts
        assert ".md" in exts

    def test_contains_image_formats(self):
        exts = get_supported_extensions()
        for ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"]:
            assert ext in exts

    def test_contains_audio_formats(self):
        exts = get_supported_extensions()
        for ext in [".mp3", ".wav", ".m4a"]:
            assert ext in exts

    def test_contains_archive_formats(self):
        exts = get_supported_extensions()
        assert ".zip" in exts

    def test_non_empty(self):
        exts = get_supported_extensions()
        assert len(exts) > 10


class TestConvertFile:
    """convert_file 函数"""

    def test_txt_file(self, tmp_path):
        fp = tmp_path / "test.txt"
        fp.write_text("Hello TXT world", encoding="utf-8")
        md, images, image_map = convert_file(fp)
        assert "Hello TXT world" in md
        assert isinstance(images, int)
        assert isinstance(image_map, list)
        assert image_map == []

    def test_csv_file(self, tmp_path):
        fp = tmp_path / "data.csv"
        fp.write_text("name,value\nfoo,1\nbar,2", encoding="utf-8")
        md, images, image_map = convert_file(fp)
        assert "foo" in md or "name" in md

    def test_md_file(self, tmp_path):
        fp = tmp_path / "notes.md"
        fp.write_text("# Hello\n\nSome content.", encoding="utf-8")
        md, images, image_map = convert_file(fp)
        assert "Hello" in md or "content" in md

    def test_nonexistent_file(self):
        """不存在的文件应抛出 FileNotFoundError"""
        with pytest.raises(FileNotFoundError):
            convert_file(Path("/nonexistent/xyz.txt"))

    def test_html_file(self, tmp_path):
        fp = tmp_path / "page.html"
        fp.write_text("<html><body><h1>Title</h1><p>Content</p></body></html>", encoding="utf-8")
        md, images, image_map = convert_file(fp)
        assert "Title" in md or "Content" in md

    def test_epub_file(self, tmp_path):
        """EPUB 文件应能处理（即使内容简单）"""
        fp = tmp_path / "book.epub"
        # Create a minimal EPUB (zip-based)
        with zipfile.ZipFile(fp, 'w') as zf:
            zf.writestr("mimetype", "application/epub+zip")
            zf.writestr("META-INF/container.xml",
                        '<?xml version="1.0"?><container><rootfiles>'
                        '<rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>'
                        '</rootfiles></container>')
            zf.writestr("OEBPS/content.opf",
                        '<?xml version="1.0"?><package version="2.0" unique-ids="uid">'
                        '<metadata><dc-metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
                        '<dc:identifier id="uid">test</dc:identifier></dc-metadata></metadata>'
                        '<manifest><item id="ch1" href="ch1.xhtml" media-type="application/xhtml+xml"/></manifest>'
                        '<spine><itemref idref="ch1"/></spine></package>')
            zf.writestr("OEBPS/ch1.xhtml",
                        '<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml">'
                        '<head><title>Test</title></head><body><p>Hello EPUB</p></body></html>')

        md, images, image_map = convert_file(fp)
        assert "EPUB" in md or "Hello" in md or len(md) > 0

    def test_empty_txt_file(self, tmp_path):
        fp = tmp_path / "empty.txt"
        fp.write_text("", encoding="utf-8")
        md, images, image_map = convert_file(fp)
        assert isinstance(md, str)


class TestConvertFileWithOutputDir:
    """convert_file 带 output_dir 参数（图片提取路径）"""

    def test_txt_with_output_dir(self, tmp_path):
        fp = tmp_path / "test.txt"
        fp.write_text("Hello", encoding="utf-8")
        out = tmp_path / "out"
        out.mkdir()
        md, images, image_map = convert_file(fp, output_dir=out)
        assert "Hello" in md
        # No images expected for txt
        assert images == 0
        assert image_map == []

    def test_with_ocr_service_none(self, tmp_path):
        """ocr_service=None 时不应崩溃"""
        fp = tmp_path / "test.txt"
        fp.write_text("Content", encoding="utf-8")
        out = tmp_path / "out"
        out.mkdir()
        md, images, image_map = convert_file(fp, output_dir=out, ocr_service=None)
        assert "Content" in md


class TestDocxImageExtraction:
    """DOCX 图片提取"""

    def _make_docx_with_image(self, tmp_path, filename="img_test.docx") -> Path:
        """Create a DOCX with a fake image embedded."""
        from docx import Document
        from docx.shared import Emu
        from docx.oxml.ns import qn, nsdecls
        from lxml import etree

        filepath = tmp_path / filename
        doc = Document()
        doc.add_paragraph("Text with image.")

        # Add a fake image reference to the document
        # We need to add an actual image file to the word/media/ part
        r = doc.add_paragraph().add_run()
        # Create a minimal PNG (1x1 pixel)
        import base64
        # Minimal 1x1 red PNG
        png_data = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8/5+hHgAHggJ/PchI7wAAAABJRU5ErkJggg=="
        )
        img_path = tmp_path / "temp_img.png"
        img_path.write_bytes(png_data)
        r.add_picture(str(img_path), width=Emu(360000))
        doc.save(str(filepath))
        return filepath

    def test_docx_image_extraction(self, tmp_path):
        fp = self._make_docx_with_image(tmp_path)
        out = tmp_path / "out"
        out.mkdir()
        md, images, image_map = convert_file(fp, output_dir=out)
        # Should extract at least one image
        assert isinstance(images, int)
        assert isinstance(image_map, list)

    def test_docx_no_images(self, tmp_path):
        """DOCX without images"""
        from docx import Document
        fp = tmp_path / "noimg.docx"
        doc = Document()
        doc.add_paragraph("No images here.")
        doc.save(str(fp))

        out = tmp_path / "out"
        out.mkdir()
        md, images, image_map = convert_file(fp, output_dir=out)
        assert images == 0
        assert image_map == []

    def test_docx_image_extraction_numeric_order(self, tmp_path):
        """DOCX 含 image1..image11 时应按数字而非字典序排列。

        字典序下：image1 < image10 < image11 < image2 → 错误
        数字序下：image1 < image2 < ... < image10 < image11 → 正确
        """
        from doc_knowledge.converters import _extract_docx_images
        import zipfile

        # 构造一个包含 11 个 media 文件名的最小 DOCX 结构
        fp = tmp_path / "many_images.docx"
        # 最小 1x1 PNG
        import base64
        png_data = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8/5+h"
            "HgAHggJ/PchI7wAAAABJRU5ErkJggg=="
        )

        with zipfile.ZipFile(fp, "w") as zf:
            zf.writestr("[Content_Types].xml", "<dummy/>")  # 满足 zip 结构即可
            # 故意按字典序的反序写入：image11 先，image1 最后
            for i in [11, 10, 9, 8, 7, 6, 5, 4, 3, 2, 1]:
                zf.writestr(f"word/media/image{i}.png", png_data)

        out = tmp_path / "out"
        out.mkdir()
        count, paths = _extract_docx_images(fp, out)
        assert count == 11

        # 验证返回顺序：image1, image2, ..., image11
        names = [p.name for p in paths]
        assert names == [f"image{i}.png" for i in range(1, 12)], (
            f"图片应按数字顺序返回，实际：{names}"
        )


class TestPptxImageExtraction:
    """PPTX 图片提取"""

    def test_pptx_without_images(self, tmp_path):
        """PPTX without images should work without error"""
        from pptx import Presentation
        fp = tmp_path / "noimg.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Hello"
        prs.save(str(fp))

        out = tmp_path / "out"
        out.mkdir()
        md, images, image_map = convert_file(fp, output_dir=out)
        assert isinstance(md, str)


class TestImageMap:
    """_build_image_map 位置匹配 + 重复引用名处理"""

    def test_image_map_preserves_duplicate_refs(self):
        """重复引用名（中文 shape 退化的 `.jpg`）应各保留独立映射条目"""
        from doc_knowledge.converters import _build_image_map

        markdown = (
            "![第一张](.jpg)\n"
            "![第二张](.jpg)\n"
            "![第三张](.jpg)\n"
        )
        paths = [
            Path("pres.pptx_images/slide1_img1.png"),
            Path("pres.pptx_images/slide1_img2.png"),
            Path("pres.pptx_images/slide2_img1.png"),
        ]
        image_map = _build_image_map(markdown, paths, "pres.pptx")
        assert isinstance(image_map, list)
        assert len(image_map) == 3
        refs = [ref for ref, _ in image_map]
        assert refs == [".jpg", ".jpg", ".jpg"], "重复引用名必须被保留"
        assert image_map[0][1] == "pres.pptx_images/slide1_img1.png"

    def test_image_map_replace_positional(self):
        """逐次替换第一个匹配项后，3 个重复 `.jpg` 引用分别指向不同路径"""
        from doc_knowledge.converters import _build_image_map

        markdown = "![a](.jpg)\n![b](.jpg)\n![c](.jpg)\n"
        paths = [
            Path("pres.pptx_images/slide1_img1.png"),
            Path("pres.pptx_images/slide1_img2.png"),
            Path("pres.pptx_images/slide2_img1.png"),
        ]
        image_map = _build_image_map(markdown, paths, "pres.pptx")

        result = markdown
        for old_name, new_path in image_map:
            result = result.replace(f"]({old_name})", f"]({new_path})", 1)

        assert "](pres.pptx_images/slide1_img1.png)" in result
        assert "](pres.pptx_images/slide1_img2.png)" in result
        assert "](pres.pptx_images/slide2_img1.png)" in result

    def test_filter_meaningless_images(self, tmp_path):
        """过滤文件过小 / 纯色的低价值图片"""
        from doc_knowledge.converters import _filter_meaningless_images
        from PIL import Image, ImageDraw

        tiny = tmp_path / "tiny.png"
        tiny.write_bytes(b"x" * 100)  # < 500B

        solid = tmp_path / "solid.bmp"
        Image.new("RGB", (100, 100), color=(255, 0, 0)).save(solid, format="BMP")

        meaningful = tmp_path / "meaningful.bmp"
        img = Image.new("RGB", (100, 100), color=(255, 255, 255))
        ImageDraw.Draw(img).rectangle([0, 0, 50, 50], fill=(0, 0, 0))
        img.save(meaningful, format="BMP")

        result = _filter_meaningless_images([tiny, solid, meaningful])
        assert meaningful in result
        assert tiny not in result
        assert solid not in result

    def test_drop_meaningless_images(self, tmp_path):
        """删除无意义图片：md 引用 + image_map 条目 + 物理文件"""
        from doc_knowledge.converters import _drop_meaningless_images
        from PIL import Image, ImageDraw

        tiny = tmp_path / "tiny.png"
        tiny.write_bytes(b"x" * 100)  # < 500B

        solid = tmp_path / "solid.bmp"
        Image.new("RGB", (100, 100), color=(255, 0, 0)).save(solid, format="BMP")

        meaningful = tmp_path / "meaningful.bmp"
        img = Image.new("RGB", (100, 100), color=(255, 255, 255))
        ImageDraw.Draw(img).rectangle([0, 0, 50, 50], fill=(0, 0, 0))
        img.save(meaningful, format="BMP")

        markdown = (
            "![tiny](.jpg)\n"
            "![solid](.jpg)\n"
            "![meaningful](.jpg)\n"
        )
        image_map = [
            (".jpg", "pres.pptx_images/slide1_img1.png"),
            (".jpg", "pres.pptx_images/slide1_img2.png"),
            (".jpg", "pres.pptx_images/slide2_img1.png"),
        ]
        image_paths = [tiny, solid, meaningful]

        new_markdown, new_map, new_paths = _drop_meaningless_images(
            markdown, image_map, image_paths
        )

        # tiny/solid：md 引用删除、image_map 条目消失、物理文件删除
        assert "![tiny]" not in new_markdown
        assert "![solid]" not in new_markdown
        assert "![meaningful]" in new_markdown
        assert len(new_map) == 1
        assert new_map[0] == (".jpg", "pres.pptx_images/slide2_img1.png")
        assert not tiny.exists()
        assert not solid.exists()
        # meaningful：引用、条目、文件均保留
        assert meaningful.exists()
        assert new_paths == [meaningful]

    def test_ocr_blockquote_injection_positional(self):
        """blockquote 应注入到各自图片引用下方，而非全部堆叠在第一张图片处

        回归：原实现用 markdown.replace("](.jpg)", "](.jpg)\\n\\n> 📷...", 1)，
        替换文本自身含 .jpg 锚点，导致后续 replace 反复匹配同一位置，
        33 张识别结果全部堆叠在第一张图片引用下。
        """
        from doc_knowledge.converters import _inject_ocr_blockquotes

        markdown = "![一](.jpg)\n甲\n\n![二](.jpg)\n乙\n\n![三](.jpg)\n丙\n"
        image_map = [
            (".jpg", "pres.pptx_images/slide1_img1.png"),
            (".jpg", "pres.pptx_images/slide1_img2.png"),
            (".jpg", "pres.pptx_images/slide2_img1.png"),
        ]
        recognized = {"slide1_img1.png": "识别甲", "slide2_img1.png": "识别丙"}

        result = _inject_ocr_blockquotes(markdown, image_map, recognized)

        # 每张已识别图片的引用下方各自追加 blockquote，互不堆叠
        assert result == (
            "![一](.jpg)\n\n> 📷 **图片识别**: 识别甲\n甲\n\n"
            "![二](.jpg)\n乙\n\n"
            "![三](.jpg)\n\n> 📷 **图片识别**: 识别丙\n丙\n"
        )
        # .jpg 引用保留（路径替换仍由调用方统一负责）
        assert result.count("](.jpg)") == 3

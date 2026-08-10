"""slide 模式（方案C：整页渲染 + 云端 VLM 三位一体意图识别）测试

覆盖：
- SlideFusionService 构造与默认参数
- _convert_pptx_to_pdf（soffice 命令）
- _render_pdf_to_pages（fitz 渲染整页 PNG）
- recognize_slides（批量送云端 VLM）
- _inject_slide_blockquotes（按页注入）
- convert_file 集成 slide 服务
"""

import base64
import json
import subprocess
import threading
from http.server import HTTPServer, BaseHTTPRequestHandler
from pathlib import Path

import pytest

from doc_knowledge.ocr.slide import SlideFusionService


class _SlideAPIHandler(BaseHTTPRequestHandler):
    """模拟 OpenAI 兼容 API：校验请求含整页图片 + slide prompt"""

    def do_POST(self):
        length = int(self.headers.get("Content-Length", 0))
        body = json.loads(self.rfile.read(length))
        msg = body["messages"]
        # [0] system, [1] user
        assert msg[0]["role"] == "system"
        user = msg[1]["content"]
        img_item = user[0]
        assert img_item["type"] == "image_url"
        assert img_item["image_url"]["url"].startswith("data:image/")
        text_item = user[1]
        assert text_item["type"] == "text"
        # 约定标记格式：prompt 必须包含标记关键词（[DK-正文] 等）
        assert "DK-正文" in text_item["text"]

        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self.end_headers()
        resp = json.dumps(
            {"choices": [{"message": {"content": (
                "[DK-标题]\n测试页\n\n[DK-概述]\n主旨\n\n"
                "[DK-结构]\n总分结构\n\n[DK-正文]\n**总**：内容"
            )}}]}
        )
        self.wfile.write(resp.encode())

    def log_message(self, format, *args):
        pass


@pytest.fixture
def mock_api_server():
    """启动模拟 API 服务器"""
    server = HTTPServer(("127.0.0.1", 0), _SlideAPIHandler)
    port = server.server_address[1]
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    yield f"http://127.0.0.1:{port}"
    server.shutdown()


def _make_minimal_pdf(path: Path, pages: int = 3):
    """用 fitz 生成最小 PDF（真实渲染目标）"""
    import fitz

    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page(width=400, height=300)
        page.insert_text((50, 80), f"Slide {i+1} test content")
    doc.save(str(path))
    doc.close()
    return path


def _make_test_image(path: Path, size=(200, 200)):
    """创建非纯色的测试图片"""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", size, color=(240, 240, 240))
    draw = ImageDraw.Draw(img)
    for i in range(0, size[0], 20):
        draw.line([(i, 0), (i, size[1])], fill=(180, 180, 180))
    draw.rectangle([30, 30, 170, 70], outline=(0, 0, 0), width=2)
    draw.text((40, 40), "Test", fill=(0, 0, 0))
    img.save(path)
    return path


class TestSlideFusionServiceInit:
    """构造与默认参数"""

    def test_defaults(self):
        svc = SlideFusionService(api_url="http://x", api_key="k")
        assert svc.model == "glm-4.6v-flash"
        assert svc.dpi == 150
        assert "DK-正文" in svc.prompt  # 约定标记：必需字段标记
        assert "[DK-结构]" in svc.prompt  # 内容逻辑结构标记
        assert "不要输出 JSON" in svc.prompt  # 标记 Markdown 输出
        # prompt 强化学（2026-08-10）：概述/结构必须输出，防止小模型（Qwen 7B）省略可选标记
        assert "必须输出，不得省略" in svc.prompt

    def test_custom_params(self):
        svc = SlideFusionService(
            api_url="http://x", api_key="k", model="gpt-4o", dpi=200,
            prompt="自定义提示", libreoffice_path="C:/soffice.exe",
        )
        assert svc.model == "gpt-4o"
        assert svc.dpi == 200
        assert svc.prompt == "自定义提示"
        assert svc.libreoffice_path == "C:/soffice.exe"


class TestConvertPptxToPdf:
    """soffice 转换命令"""

    def test_uses_soffice_headless(self, tmp_path, monkeypatch):
        svc = SlideFusionService(api_url="http://x", api_key="k")
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(b"PK")

        monkeypatch.setattr(
            "doc_knowledge.ocr.slide._find_soffice", lambda: "C:/fake/soffice.exe"
        )
        calls = {}

        def fake_run(cmd, **kwargs):
            calls["cmd"] = cmd
            calls["kwargs"] = kwargs
            # 模拟 soffice 产生 PDF 到 --outdir
            i = cmd.index("--outdir")
            outdir = Path(cmd[i + 1])
            pdf = outdir / f"{Path(cmd[-1]).stem}.pdf"
            pdf.write_bytes(b"%PDF")
            return subprocess.CompletedProcess(cmd, 0)

        monkeypatch.setattr(subprocess, "run", fake_run)
        out = tmp_path / "work"
        out.mkdir()
        pdf = svc._convert_pptx_to_pdf(pptx, out)

        assert pdf == out / "test.pdf"
        cmd = calls["cmd"]
        assert "--headless" in cmd
        assert "--convert-to" in cmd
        assert "pdf" in cmd
        assert str(pptx) in cmd

    def test_missing_soffice_raises(self, tmp_path, monkeypatch):
        svc = SlideFusionService(api_url="http://x", api_key="k")
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(b"PK")

        monkeypatch.setattr("doc_knowledge.ocr.slide._find_soffice", lambda: "")
        with pytest.raises(RuntimeError, match="LibreOffice"):
            svc._convert_pptx_to_pdf(pptx, tmp_path)


class TestRenderPdfToPages:
    """fitz 渲染整页 PNG"""

    def test_renders_all_pages(self, tmp_path):
        svc = SlideFusionService(api_url="http://x", api_key="k", dpi=100)
        pdf = _make_minimal_pdf(tmp_path / "src.pdf", pages=3)
        pages = svc._render_pdf_to_pages(pdf, tmp_path)

        assert len(pages) == 3
        for i, p in enumerate(pages):
            assert p.exists()
            assert p.name == f"page{i + 1}.png"
            # 非空图片
            assert p.stat().st_size > 0

    def test_png_size_scales_with_dpi(self, tmp_path):
        low = SlideFusionService(api_url="http://x", api_key="k", dpi=72)
        high = SlideFusionService(api_url="http://x", api_key="k", dpi=200)
        pdf = _make_minimal_pdf(tmp_path / "src.pdf", pages=1)

        p1 = low._render_pdf_to_pages(pdf, tmp_path / "d1")[0]
        p2 = high._render_pdf_to_pages(pdf, tmp_path / "d2")[0]
        assert p2.stat().st_size > p1.stat().st_size


class TestRecognizeSlides:
    """整页 PNG 批量送云端 VLM"""

    def test_returns_page_indexed_results(self, mock_api_server, tmp_path):
        svc = SlideFusionService(api_url=mock_api_server, api_key="k")
        pages = [
            _make_test_image(tmp_path / "page1.png"),
            _make_test_image(tmp_path / "page2.png"),
        ]
        results = svc.recognize_slides(pages)

        assert len(results) == 2
        assert 1 in results and 2 in results
        for num, text in results.items():
            assert "[DK-正文]" in text  # 约定标记：必需字段

    def test_429_retry_with_backoff(self, tmp_path, monkeypatch):
        """429 限流后按退避重试，最终成功

        回归：智谱免费 VLM（glm-4.6v-flash）连续请求触发 HTTP 429，
        需要指数退避等待额度恢复后重试。
        """
        svc = SlideFusionService(api_url="http://x", api_key="k")
        page = _make_test_image(tmp_path / "page1.png")
        calls = {"n": 0}

        def fake_recognize(img):
            calls["n"] += 1
            if calls["n"] == 1:
                return "[图片识别失败: HTTP Error 429: Too Many Requests]"
            if calls["n"] == 2:
                return "[图片识别失败: HTTP Error 429: Too Many Requests]"
            return "[DK-正文]\n限流恢复后成功"

        monkeypatch.setattr(svc._vision, "recognize_image", fake_recognize)
        monkeypatch.setattr("doc_knowledge.ocr.slide.time.sleep", lambda s: None)

        results = svc.recognize_slides([page], max_retries=3, retry_base_delay=1.0)
        assert calls["n"] == 3
        assert results == {1: "[DK-正文]\n限流恢复后成功"}

    def test_retry_exhausted_returns_last_error(self, tmp_path, monkeypatch):
        """重试耗尽后返回最后一次错误（不抛异常）"""
        svc = SlideFusionService(api_url="http://x", api_key="k")
        page = _make_test_image(tmp_path / "page1.png")

        def fake_recognize(img):
            return "[图片识别失败: HTTP Error 429: Too Many Requests]"

        monkeypatch.setattr(svc._vision, "recognize_image", fake_recognize)
        monkeypatch.setattr("doc_knowledge.ocr.slide.time.sleep", lambda s: None)

        results = svc.recognize_slides([page], max_retries=2, retry_base_delay=1.0)
        assert "429" in results[1]

    def test_empty_response_triggers_retry(self, tmp_path, monkeypatch):
        """空串响应应视为失败触发重试，而非静默成功

        回归：真实 75 页 PPT 验证中发现第 30 页 VLM 返回空字符串，
        _recognize_with_retry 将空串误判为成功，导致整页注入静默丢失。
        """
        svc = SlideFusionService(api_url="http://x", api_key="k")
        page = _make_test_image(tmp_path / "page1.png")
        calls = {"n": 0}

        def fake_recognize(img):
            calls["n"] += 1
            if calls["n"] == 1:
                return ""  # VLM 返回空串（非 [错误标记] 形态）
            return "[DK-正文]\n空串重试后成功"

        monkeypatch.setattr(svc._vision, "recognize_image", fake_recognize)
        monkeypatch.setattr("doc_knowledge.ocr.slide.time.sleep", lambda s: None)

        results = svc.recognize_slides([page], max_retries=3, retry_base_delay=1.0)
        assert calls["n"] == 2  # 空串触发重试，第二次成功
        assert results == {1: "[DK-正文]\n空串重试后成功"}

    def test_whitespace_response_triggers_retry(self, tmp_path, monkeypatch):
        """纯空白响应同样视为失败触发重试"""
        svc = SlideFusionService(api_url="http://x", api_key="k")
        page = _make_test_image(tmp_path / "page1.png")
        calls = {"n": 0}

        def fake_recognize(img):
            calls["n"] += 1
            if calls["n"] == 1:
                return "   \n  "  # 纯空白
            return "[DK-正文]\n空白重试后成功"

        monkeypatch.setattr(svc._vision, "recognize_image", fake_recognize)
        monkeypatch.setattr("doc_knowledge.ocr.slide.time.sleep", lambda s: None)

        results = svc.recognize_slides([page], max_retries=3, retry_base_delay=1.0)
        assert calls["n"] == 2
        assert results == {1: "[DK-正文]\n空白重试后成功"}

    def test_auto_retry_failed_pass(self, tmp_path, monkeypatch):
        """自动二次补跑：首轮单页重试也失败 → 补跑第二轮成功（默认开启）

        限流随时间恢复：全页识别完后再隔一轮，失败页补跑成功率显著提升。
        设计区分两层：首轮 _recognize_with_retry（max_retries=3 → 4 次尝试）
        全失败，只有补跑第二轮（第 5 次调用）才成功。
        """
        svc = SlideFusionService(api_url="http://x", api_key="k")
        p1 = _make_test_image(tmp_path / "page1.png")
        p2 = _make_test_image(tmp_path / "page2.png")
        counts = {}

        def fake_recognize(img):
            counts[img.name] = counts.get(img.name, 0) + 1
            n = counts[img.name]
            # page1：首轮 4 次尝试（max_retries=3）全失败，补跑第 5 次成功
            if img.name == "page1.png" and n <= 4:
                return "[图片识别失败: HTTP Error 429: Too Many Requests]"
            if img.name == "page1.png":
                return "[DK-正文]\n第一页补跑成功"
            return "[DK-正文]\n第二页"

        monkeypatch.setattr(svc._vision, "recognize_image", fake_recognize)
        monkeypatch.setattr("doc_knowledge.ocr.slide.time.sleep", lambda s: None)

        results = svc.recognize_slides([p1, p2], max_retries=3, retry_base_delay=1.0)
        assert counts["page1.png"] == 5  # 首轮 4 次失败 + 自动补跑 1 次成功
        assert counts["page2.png"] == 1  # 成功页不补跑
        assert "补跑成功" in results[1]

    def test_no_auto_retry_when_disabled(self, tmp_path, monkeypatch):
        """retry_failed_pass=False 时不做二次补跑，失败页保留错误结果"""
        svc = SlideFusionService(api_url="http://x", api_key="k")
        p1 = _make_test_image(tmp_path / "page1.png")
        counts = {"n": 0}

        def fake_recognize(img):
            counts["n"] += 1
            return "[图片识别失败: HTTP Error 429: Too Many Requests]"

        monkeypatch.setattr(svc._vision, "recognize_image", fake_recognize)
        monkeypatch.setattr("doc_knowledge.ocr.slide.time.sleep", lambda s: None)

        results = svc.recognize_slides(
            [p1], retry_failed_pass=False, max_retries=3, retry_base_delay=1.0,
        )
        assert counts["n"] == 4  # 仅单页 4 次尝试（max_retries=3），无补跑轮
        assert "429" in results[1]


class TestProcessPptxFullFlow:
    """soffice → pdf → 整页 PNG → VLM 完整编排"""

    def test_process_pptx_orchestrates(self, mock_api_server, tmp_path, monkeypatch):
        svc = SlideFusionService(api_url=mock_api_server, api_key="k")
        pptx = tmp_path / "test.pptx"
        pptx.write_bytes(b"PK")

        # 编排内部步骤：PDF 用真实渲染，soffice 用假命令
        def fake_run(cmd, **kwargs):
            # 从 soffice 命令提取 --outdir 与输入文件 stem
            i = cmd.index("--outdir")
            outdir = Path(cmd[i + 1])
            stem = Path(cmd[-1]).stem
            _make_minimal_pdf(outdir / f"{stem}.pdf", pages=2)
            return subprocess.CompletedProcess(cmd, 0)

        monkeypatch.setattr(
            "doc_knowledge.ocr.slide._find_soffice", lambda: "C:/fake/soffice.exe"
        )
        monkeypatch.setattr(subprocess, "run", fake_run)

        results = svc.process_pptx(pptx, tmp_path, verbose=False)
        assert len(results) == 2
        assert 1 in results and 2 in results


class TestInjectSlideBlockquotes:
    """按 <!-- Slide number: N --> 边界注入整页理解"""

    def test_injects_per_slide(self):
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = (
            "<!-- Slide number: 1 -->\n\n# 标题\n\n"
            "<!-- Slide number: 2 -->\n\n# 标题2\n\n"
        )
        results = {1: "第一页理解", 2: "第二页理解"}
        out = _inject_slide_blockquotes(md, results)

        assert out.count("整页理解") == 2
        assert "整页理解" in out.split("<!-- Slide number: 2 -->")[0]
        assert "第一页理解" in out
        assert "第二页理解" in out

    def test_missing_pages_left_untouched(self):
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = "<!-- Slide number: 1 -->\n\nA\n\n<!-- Slide number: 2 -->\n\nB\n"
        out = _inject_slide_blockquotes(md, {1: "仅第一页"})
        assert "仅第一页" in out
        # 第 2 页无结果，保持原样
        assert out.count("整页理解") == 1

    def test_failed_pages_inject_fallback_note(self):
        """失败页（错误占位符）注入降级提示，而非暴露错误堆栈

        回归：真实 75 页 PPT 中 16 页限流失败，原实现把
        `[图片识别失败: HTTP Error 429...]` 原文注入 markdown，
        观感差且暴露技术细节。现降级为可读中文提示。
        """
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = "<!-- Slide number: 1 -->\n\nA\n\n<!-- Slide number: 2 -->\n\nB\n"
        results = {
            1: '{"page_summary": "成功"}',
            2: "[图片识别失败: HTTP Error 429: Too Many Requests]",
        }
        out = _inject_slide_blockquotes(md, results)

        assert '{"page_summary": "成功"}' in out
        assert "HTTP Error 429" not in out       # 不暴露错误堆栈
        assert "识别失败" in out                   # 降级提示
        assert "原始文字已保留" in out


class TestConvertFileIntegration:
    """convert_file 使用 SlideFusionService 处理 PPTX"""

    def test_slide_service_injects_into_markdown(self, tmp_path, monkeypatch):
        from docx import Document
        from doc_knowledge.converters import convert_file

        # 真实 PPTX：markitdown 会输出 <!-- Slide number: 1 --> 标记
        from pptx import Presentation

        pptx = tmp_path / "pres.pptx"
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "市场规模"
        slide1.placeholders[1].text = "2026年超700亿元"
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])
        slide2.shapes.title.text = "竞争格局"
        slide2.placeholders[1].text = "四大阵营"
        prs.save(str(pptx))

        out = tmp_path / "out"
        out.mkdir()

        svc = SlideFusionService(api_url="http://x", api_key="k")
        monkeypatch.setattr(
            svc, "process_pptx",
            lambda pptx, out, verbose=False: {1: "整页：本页讲市场规模", 2: "整页：本页讲竞争"},
        )

        md, images, image_map = convert_file(pptx, output_dir=out, ocr_service=svc)

        assert "整页：本页讲市场规模" in md
        assert "整页：本页讲竞争" in md
        assert "整页理解" in md
        assert "<!-- Slide number: 1 -->" in md

    def test_regular_ocr_service_not_affected(self, tmp_path):
        """OCRService（图像级）不应触发 slide 注入"""
        from doc_knowledge.ocr.cloud import CloudOCRService

        svc = CloudOCRService(api_url="http://x", api_key="k")
        assert not isinstance(svc, SlideFusionService)

    def test_slide_mode_skips_embedded_images(self, tmp_path, monkeypatch):
        """slide 模式不产内嵌图：images=0、image_map=[]、md 无图片引用、无 _images/ 残留

        回归：slide 模式识别对象是整页截图而非内嵌图，内嵌图是噪音残留。
        """
        from PIL import Image, ImageDraw
        from pptx import Presentation

        from doc_knowledge.converters import convert_file

        # 真实 PPTX：含标题 + 1 张内嵌图片（MarkItDown 会输出图片引用）
        pptx = tmp_path / "pres.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "市场规模"
        img_path = tmp_path / "tmp.png"
        # 多色复杂图：确保不被 ImageFilter 判为无意义（体积>500B、分辨率够、非纯色），
        # 使 MarkItDown 的图片引用在正常管道中存活，才能验证 slide 分支主动删除。
        img = Image.new("RGB", (400, 300), (255, 255, 255))
        d = ImageDraw.Draw(img)
        d.rectangle([10, 10, 120, 90], fill=(200, 30, 30))
        d.rectangle([140, 40, 260, 140], fill=(30, 60, 200))
        d.rectangle([60, 160, 200, 260], fill=(30, 160, 60))
        d.rectangle([220, 180, 360, 280], fill=(0, 0, 0))
        img.save(img_path, format="PNG")
        slide.shapes.add_picture(str(img_path), 0, 0)
        prs.save(str(pptx))

        out = tmp_path / "out"
        out.mkdir()

        svc = SlideFusionService(api_url="http://x", api_key="k")
        monkeypatch.setattr(
            svc, "process_pptx",
            lambda pptx, out, verbose=False: {1: "整页：本页讲市场规模"},
        )

        md, images, image_map = convert_file(pptx, output_dir=out, ocr_service=svc)

        assert "整页理解" in md
        assert "![" not in md                    # 内嵌图引用被删除
        assert images == 0                        # 不提取内嵌图
        assert image_map == []                    # 不建图片映射
        assert not (out / "pres.pptx_images").exists()


class TestSlideStructuredOutput:
    """结构化输出：_parse_slide_result / _render_slide_page / _inject_slide_blockquotes"""

    def test_parse_slide_result_marked_markdown(self):
        """标准约定标记输入 → 返回正确 dict（title/overview/structure/body）"""
        from doc_knowledge.converters import _parse_slide_result

        text = (
            "[DK-标题]\n市场规模\n\n[DK-概述]\n2026超700亿\n\n"
            "[DK-结构]\n总分结构\n\n[DK-正文]\n**总**：x\n- a\n- b"
        )
        result = _parse_slide_result(text)
        assert result is not None
        assert result["title"] == "市场规模"
        assert result["overview"] == "2026超700亿"
        assert result["structure"] == "总分结构"
        assert result["body"] == "**总**：x\n- a\n- b"

    def test_parse_slide_result_only_body(self):
        """仅 [DK-正文] → {"body": ...}（title/overview/structure 缺省不出现）"""
        from doc_knowledge.converters import _parse_slide_result

        result = _parse_slide_result("[DK-正文]\n仅正文内容")
        assert result == {"body": "仅正文内容"}

    def test_parse_slide_result_missing_body(self):
        """无 [DK-正文] 标记 → None"""
        from doc_knowledge.converters import _parse_slide_result

        assert _parse_slide_result("[DK-标题]\n只有标题") is None
        assert _parse_slide_result("整页：普通文本") is None
        assert _parse_slide_result('{"page_summary": "x"}') is None   # 旧 JSON schema 缺 body

    def test_parse_slide_result_empty_body(self):
        """[DK-正文] 后为空 → None（走降级）"""
        from doc_knowledge.converters import _parse_slide_result

        assert _parse_slide_result("[DK-标题]\nx\n\n[DK-正文]\n   ") is None
        assert _parse_slide_result("[DK-正文]\n") is None

    def test_parse_slide_result_invalid(self):
        """失败占位符 → None"""
        from doc_knowledge.converters import _parse_slide_result

        assert _parse_slide_result("[图片识别失败: HTTP 429]") is None  # 错误占位符
        assert _parse_slide_result("[图片解析失败: timeout]") is None

    def test_parse_slide_result_multiline_body(self):
        """标记 Markdown：body 内真实换行/制表符原生成立，无需转义

        回归：JSON 时代 GLM 的 body 值内真实换行/制表符导致 json.loads 抛
        Invalid control character → 整页降级。标记 Markdown 无转义负担，
        多行正文直接解析成功。
        """
        from doc_knowledge.converters import _parse_slide_result

        text = (
            "[DK-标题]\n风险\n\n[DK-概述]\no\n\n[DK-结构]\n并列\n\n"
            "[DK-正文]\n**提示**：谨慎\n\t- 风险一\n- 风险二"
        )
        result = _parse_slide_result(text)
        assert result is not None
        assert result["body"] == "**提示**：谨慎\n\t- 风险一\n- 风险二"

    def test_parse_slide_result_leading_text(self):
        """标记前有额外解释文本（未严格守格式）仍能解析——正则提取不受前置内容影响"""
        from doc_knowledge.converters import _parse_slide_result

        text = "识别结果如下：\n[DK-标题]\n市场规模\n\n[DK-正文]\n核心结论"
        result = _parse_slide_result(text)
        assert result is not None
        assert result["title"] == "市场规模"
        assert result["body"] == "核心结论"

    def test_inject_slide_multiline_structured(self):
        """标记 Markdown 多行 body 注入后走结构化渲染，产物不出现裸 JSON 代码块

        回归：JSON 时代 GLM 返回降级为 `> 📊 **整页理解**: ```json … ``` `，
        json.loads 对产物解析 89% 非法；标记迁移后应渲染为可读 Markdown。
        """
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = "<!-- Slide number: 1 -->\n\n原文\n"
        results = {
            1: "[DK-标题]\n风险\n\n[DK-概述]\no\n\n[DK-结构]\n并列\n\n"
               "[DK-正文]\n**提示**：谨慎\n\t- 风险一\n- 风险二",
        }
        out = _inject_slide_blockquotes(md, results)

        assert "📊 **风险**（并列）" in out
        assert "**提示**：谨慎" in out
        assert "- 风险一" in out
        assert "- 风险二" in out
        assert "```json" not in out          # 不出现裸 JSON 代码块
        assert "整页理解" not in out          # 不降级为旧 blockquote

    def test_render_slide_page_structured(self):
        """总分结构：标题（结构）+ 概述 + body + 原文折叠"""
        from doc_knowledge.converters import _render_slide_page

        page_text = "市场规模预测\n2026年市场超700亿元"
        result = {
            "title": "市场规模预测",
            "overview": "2026年 AI 市场超700亿元",
            "structure": "总分结构",
            "body": "**核心结论**：市场持续增长\n\n| 年份 | 规模 |\n|------|------|\n| 2024 | 280 |",
        }
        out = _render_slide_page(page_text, result)

        assert "📊 **市场规模预测**（总分结构）" in out
        assert "**概述**：2026年 AI 市场超700亿元" in out
        assert "| 年份 | 规模 |" in out                    # body 表格保留
        assert "markitdown 原文（兜底参考）" in out          # 原文折叠
        assert "2026年市场超700亿元" in out                  # 原文内容保留

    def test_inject_slide_blockquotes_structured_multi_page(self):
        """多页结构化注入：不串页，原文各归各页"""
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = (
            "<!-- Slide number: 1 -->\n\n第一页原文A\n\n"
            "<!-- Slide number: 2 -->\n\n第二页原文B\n"
        )
        results = {
            1: "[DK-标题]\n页一\n\n[DK-概述]\no1\n\n[DK-结构]\n并列\n\n[DK-正文]\nb1",
            2: "[DK-标题]\n页二\n\n[DK-概述]\no2\n\n[DK-结构]\n递进\n\n[DK-正文]\nb2",
        }
        out = _inject_slide_blockquotes(md, results)

        assert "📊 **页一**（并列）" in out
        assert "📊 **页二**（递进）" in out
        assert "第一页原文A" in out
        assert "第二页原文B" in out
        assert out.count("<details>") == 2                   # 每页一个原文折叠
        # 页1 折叠含 A 不含 B（原文不串页）
        page1_detail = out[out.index("页一"): out.index("页二")]
        assert "第一页原文A" in page1_detail

    def test_inject_slide_blockquotes_json_fallback(self):
        """非 JSON 文本 → 降级旧 blockquote（保留文本）"""
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = "<!-- Slide number: 1 -->\n\n原文\n"
        out = _inject_slide_blockquotes(md, {1: "整页：普通文本"})

        assert "> 📊 **整页理解**: 整页：普通文本" in out
        assert "原文" in out

    def test_inject_slide_blockquotes_failure_fallback(self):
        """失败占位符 → 降级提示，不暴露错误堆栈"""
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = "<!-- Slide number: 1 -->\n\n原文\n"
        out = _inject_slide_blockquotes(md, {1: "[图片识别失败: HTTP Error 429]"})

        assert "识别失败" in out
        assert "HTTP Error 429" not in out
        assert "原文" in out

    def test_render_slide_page_body_as_dict(self):
        """VLM 把 body 输出为嵌套 dict → 渲染为 Markdown 分节+列表，避免裸 JSON

        回归：真实 10 页验证对章节小结页**稳定复现**（两次均为 dict），
        直接 .strip() 崩溃。dict 应转为 `### 键` + `- 值` 的 Markdown 形态。
        """
        from doc_knowledge.converters import _render_slide_page

        page_text = "市场规模预测"
        result = {
            "title": "市场规模预测",
            "overview": "2026超700亿",
            "structure": "总分结构",
            "body": {"核心结论": "市场持续增长", "支撑": ["用户增长", "政策支持"]},
        }
        out = _render_slide_page(page_text, result)

        assert "📊 **市场规模预测**（总分结构）" in out
        assert "**概述**：2026超700亿" in out
        assert "### 核心结论" in out        # dict 键 → 小标题
        assert "- 市场持续增长" in out      # dict 值（str）→ 列表项
        assert "### 支撑" in out
        assert "- 用户增长" in out          # dict 值（list）→ 列表
        assert "- 政策支持" in out
        assert '{"核心结论"' not in out     # 不出现裸 JSON

    def test_render_slide_page_body_as_list(self):
        """body 直接为 list → 渲染为项目符号列表"""
        from doc_knowledge.converters import _render_slide_page

        page_text = "要点"
        result = {
            "title": "要点",
            "overview": "o",
            "structure": "并列",
            "body": ["第一点", "第二点", "第三点"],
        }
        out = _render_slide_page(page_text, result)

        assert "- 第一点" in out
        assert "- 第二点" in out
        assert "- 第三点" in out
        assert '["第一点"' not in out       # 不出现裸 JSON

    def test_inject_no_leading_blank_line(self):
        """markdown 以 <!-- Slide number: 1 --> 开头时，输出不以空行开头

        回归：真实产物第一行异常空行——re.split 产生的 parts[0] 为空串，
        '\n'.join 在开头拼出一个 '\n'。注入不应改变 markdown 头部格式。
        """
        from doc_knowledge.converters import _inject_slide_blockquotes

        md = "<!-- Slide number: 1 -->\n\n第一页原文\n"
        out = _inject_slide_blockquotes(
            md, {1: "[DK-标题]\n页一\n\n[DK-概述]\no\n\n[DK-结构]\n并列\n\n[DK-正文]\nb"}
        )

        assert not out.startswith("\n")
        assert out.startswith("<!-- Slide number: 1 -->")

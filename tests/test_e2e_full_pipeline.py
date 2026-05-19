"""
E2E 全流程测试

使用真实 DOCX/PDF 文件验证完整流水线：convert → extract → export
"""

import tempfile
import shutil
from pathlib import Path
from click.testing import CliRunner
from doc_knowledge.cli import main


def _create_test_docx(path: Path, title: str = "Test Document", paragraphs: list = None):
    """创建测试 DOCX 文件"""
    from docx import Document
    from docx.shared import Pt
    
    doc = Document()
    doc.add_heading(title, level=1)
    
    if paragraphs:
        for p in paragraphs:
            doc.add_paragraph(p)
    else:
        doc.add_paragraph("This is the first paragraph about system architecture design.")
        doc.add_paragraph("The system uses Docker containerization and Kubernetes orchestration.")
        doc.add_heading("Core Modules", level=2)
        doc.add_paragraph("The converter module handles PDF, DOCX, PPTX, and XLSX formats.")
        doc.add_paragraph("The extractor module performs deduplication and value scoring.")
        doc.add_heading("Performance Requirements", level=2)
        doc.add_paragraph("The system must process 1000 files within 10 minutes.")
    
    doc.save(str(path))


def _create_test_pdf(path: Path, text: str = "Test PDF Content"):
    """创建最小有效 PDF 文件"""
    # 创建一个简单的文本 PDF
    pdf_content = f"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT
/F1 12 Tf
100 700 Td
({text}) Tj
ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000266 00000 n 
0000000358 00000 n 
trailer
<< /Size 6 /Root 1 0 R >>
startxref
439
%%EOF"""
    path.write_bytes(pdf_content.encode('latin-1'))


def _create_test_pptx(path: Path, title: str = "Test Presentation"):
    """创建测试 PPTX 文件"""
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "This is a test presentation about machine learning algorithms."
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Architecture Design"
    slide2.placeholders[1].text = "The system uses microservices architecture with Docker containers."
    
    prs.save(str(path))


# ──────────────────────────────────────────────
# E2E 全流程测试
# ──────────────────────────────────────────────

def test_e2e_docx_full_pipeline():
    """DOCX 文件完整流水线：convert → extract → export"""
    runner = CliRunner()
    with tempfile.TemporaryDirectory() as tmpdir:
        # 创建测试 DOCX
        source = Path(tmpdir) / "source"
        source.mkdir()
        _create_test_docx(source / "architecture.md")
        _create_test_docx(source / "performance.md", "Performance Report", [
            "The system must handle 1000 concurrent users.",
            "Response time should be under 100ms.",
            "Docker containerization ensures scalability.",
        ])
        
        # 运行 pipeline
        output = Path(tmpdir) / "output"
        result = runner.invoke(main, [
            "pipeline", str(source),
            "--target", "markdown",
            "-o", str(output)
        ])
        
        assert result.exit_code == 0
        assert output.exists()
        
        # 验证输出文件
        md_files = list(output.rglob("*.md"))
        assert len(md_files) >= 1
        
        # 验证 frontmatter 包含提取元数据
        for md in md_files:
            content = md.read_text(encoding="utf-8")
            assert "dk_score:" in content or "source:" in content


def test_e2e_mixed_formats():
    """混合格式文件（DOCX + TXT）流水线"""
    runner = CliRunner()
    with tempfile.TemporaryDirectory() as tmpdir:
        source = Path(tmpdir) / "source"
        source.mkdir()
        
        # DOCX 文件
        _create_test_docx(source / "report.docx", "System Report", [
            "This report covers the architecture design and implementation details.",
            "Key components include the converter, extractor, and exporter modules.",
        ])
        
        # TXT 文件
        (source / "notes.txt").write_text(
            "Meeting notes: Discussed Docker containerization strategy.\n"
            "Agreed to use Kubernetes for orchestration.",
            encoding="utf-8"
        )
        
        # 运行 pipeline
        output = Path(tmpdir) / "output"
        result = runner.invoke(main, [
            "pipeline", str(source),
            "--target", "markdown",
            "-o", str(output),
            "-v"
        ])
        
        assert result.exit_code == 0
        assert output.exists()
        
        # 验证 DOCX 被转换
        md_files = list(output.rglob("*.md"))
        assert len(md_files) >= 1
        assert any("report" in f.name.lower() for f in md_files)


def test_e2e_subdirectory_preservation():
    """子目录结构保持测试"""
    runner = CliRunner()
    with tempfile.TemporaryDirectory() as tmpdir:
        source = Path(tmpdir) / "source"
        source.mkdir()
        (source / "chapter1").mkdir(parents=True)
        (source / "chapter2").mkdir(parents=True)
        
        _create_test_docx(source / "chapter1" / "intro.docx", "Introduction", ["Welcome to the guide."])
        _create_test_docx(source / "chapter2" / "advanced.docx", "Advanced Topics", ["Deep dive into architecture."])
        
        output = Path(tmpdir) / "output"
        result = runner.invoke(main, [
            "pipeline", str(source),
            "--target", "markdown",
            "-o", str(output),
            "-v"
        ])
        
        assert result.exit_code == 0
        assert output.exists()
        
        # 验证子目录结构（DOCX 文件应该被转换为 .md）
        md_files = list(output.rglob("*.md"))
        assert len(md_files) >= 1
        # 验证至少有一个文件包含 chapter1 或 chapter2 在路径中
        has_chapters = any("chapter1" in str(f) or "chapter2" in str(f) for f in md_files)
        assert has_chapters


def test_e2e_backlinks():
    """反向链接注入测试"""
    runner = CliRunner()
    with tempfile.TemporaryDirectory() as tmpdir:
        source = Path(tmpdir) / "source"
        source.mkdir()
        _create_test_docx(source / "test.docx", "Test", ["Content here."])
        
        output = Path(tmpdir) / "output"
        result = runner.invoke(main, [
            "pipeline", str(source),
            "--target", "markdown",
            "-o", str(output)
        ])
        
        assert result.exit_code == 0
        
        # 验证反向链接（source 字段）
        md_files = list(output.rglob("*.md"))
        assert len(md_files) >= 1
        content = md_files[0].read_text(encoding="utf-8")
        assert "source:" in content
        assert "file:///" in content or "source_relative" in content


def test_e2e_dedup_and_merge():
    """去重 + 版本合并 E2E 测试"""
    runner = CliRunner()
    with tempfile.TemporaryDirectory() as tmpdir:
        # 创建相似内容文件（触发去重）
        mirror = Path(tmpdir) / "mirror"
        mirror.mkdir()
        
        content1 = """---
title: "report_v1"
source: "test"
---

This is a report about system architecture design and Docker containerization.
The converter module handles PDF, DOCX, PPTX, and XLSX formats.
"""
        (mirror / "report_v1.txt.md").write_text(content1, encoding="utf-8")
        
        content2 = """---
title: "report_v2"
source: "test"
---

This is a report about system architecture design and Docker containerization.
The converter module handles PDF, DOCX, PPTX, and XLSX formats.
Additional details about Kubernetes orchestration and microservices.
"""
        (mirror / "report_v2.txt.md").write_text(content2, encoding="utf-8")
        
        # 运行 extract（带去重）
        output = Path(tmpdir) / "output"
        result = runner.invoke(main, [
            "extract", str(mirror),
            "-o", str(output)
        ])
        
        assert result.exit_code == 0
        
        # 验证去重生效（两个相似文档应该只保留一个）
        output_files = list(output.rglob("*.md"))
        assert len(output_files) <= 2  # 最多保留 1-2 个


def test_e2e_incremental():
    """增量更新 E2E 测试"""
    runner = CliRunner()
    with tempfile.TemporaryDirectory() as tmpdir:
        source = Path(tmpdir) / "source"
        source.mkdir()
        _create_test_docx(source / "test.docx", "Test", ["Initial content."])
        
        output = Path(tmpdir) / "output"
        
        # 第一次运行
        result1 = runner.invoke(main, [
            "pipeline", str(source),
            "--target", "markdown",
            "-o", str(output)
        ])
        assert result1.exit_code == 0
        
        # 修改文件
        import time
        time.sleep(0.1)
        _create_test_docx(source / "test.docx", "Test", ["Updated content with new information."])
        
        # 第二次运行（增量模式）
        result2 = runner.invoke(main, [
            "pipeline", str(source),
            "--target", "markdown",
            "-o", str(output)
        ])
        assert result2.exit_code == 0

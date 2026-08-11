"""文档转换器包 - 基于 MarkItDown 封装"""

import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


_PDF_TEXT_DENSITY_THRESHOLD = 50  # 平均每页字符数阈值


def _pdf_has_text_layer(pdf_path: Path) -> bool:
    """检测 PDF 是否有可提取的文字层（按平均每页字符密度判断）

    旧策略"任意一页 >10 字符即认为有文字"会把扫描件中的页眉/水印误判为文字层。
    新策略：总文字字符数 / 页数 ≥ 50 才认为是真正的文字 PDF。

    阈值依据：典型水印（"Page 1 of 3"、"Confidential"）通常 <30 字符/页；
    正常段落文字通常 >100 字符/页。50 留出足够缓冲。

    导入或读取失败时返回 True（不确定时倾向跳过 OCR，避免无效 API 调用）。
    """
    try:
        import fitz
    except ImportError:
        return True

    try:
        doc = fitz.open(str(pdf_path))
        page_count = len(doc)
        if page_count == 0:
            doc.close()
            return False
        total_chars = sum(len(page.get_text().strip()) for page in doc)
        doc.close()
        return (total_chars / page_count) >= _PDF_TEXT_DENSITY_THRESHOLD
    except Exception as e:
        logger.debug("PDF 文字层检测失败，按有文字层处理: %s (%s)", pdf_path.name, e)
        return True


def _render_pdf_pages(pdf_path: Path, output_dir: Path) -> list[Path]:
    """将 PDF 每一页渲染为 PNG 图片，返回图片路径列表"""
    try:
        import fitz
    except ImportError:
        return []

    images = []
    try:
        doc = fitz.open(str(pdf_path))
        images_dir = output_dir / f"{pdf_path.name}_images"
        images_dir.mkdir(parents=True, exist_ok=True)

        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            img_path = images_dir / f"page{i + 1}.png"
            pix.save(str(img_path))
            images.append(img_path)

        doc.close()
    except Exception as e:
        logger.warning("PDF 页面渲染失败: %s (%s)", pdf_path.name, e)

    return images


def convert_file(
    filepath: Path,
    output_dir: Optional[Path] = None,
    ocr_service=None,
    verbose: bool = False,
) -> tuple[str, int, list[tuple[str, str]]]:
    """
    使用 MarkItDown 将文件转换为 Markdown

    Args:
        filepath: 源文件路径
        output_dir: 输出目录（用于保存图片）
        ocr_service: 可选的 OCR 服务（统一接口，详见 ocr/base.py）
        verbose: 是否输出详细信息

    Returns:
        (Markdown 文本, 提取的图片数量, 图片映射 [(ref_name, new_path), ...])
    """
    if not filepath.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")

    from markitdown import MarkItDown
    md = MarkItDown()
    result = md.convert(str(filepath))
    markdown = result.text_content

    images_extracted = 0
    image_paths: list[Path] = []

    # 图片型 PDF 处理：MarkItDown 返回空内容时，渲染页面 + OCR
    if filepath.suffix.lower() == '.pdf' and output_dir is not None:
        content_text = markdown.strip()
        # 去掉 frontmatter 等元数据后判断是否为空
        if not content_text or len(content_text) < 50:
            if not _pdf_has_text_layer(filepath):
                pages = _render_pdf_pages(filepath, output_dir)
                if pages:
                    if ocr_service:
                        results = ocr_service.recognize_batch(pages, verbose=verbose)
                        lines = [f"# {filepath.stem}\n"]
                        for i, page_img in enumerate(pages):
                            text = results.get(page_img, "[识别失败]")
                            lines.append(f"## 第 {i + 1} 页\n\n{text}\n")
                        markdown = "\n".join(lines)
                        images_extracted += len(pages)
                    else:
                        lines = [
                            f"# {filepath.stem}\n",
                            "> ⚠️ **图片型 PDF**：此文档为扫描件或图片型 PDF，"
                            "当前未启用 OCR 功能。\n",
                            "> 启用方法：`doc-knowledge convert <dir> --ocr cloud`"
                            " 或配置 `~/.doc-knowledge/config.yaml`\n",
                        ]
                        markdown = "\n".join(lines)
                        images_extracted += len(pages)
                        image_paths.extend(pages)

    from doc_knowledge.ocr.slide import SlideFusionService

    # slide 模式（方案C）：整页渲染 + 云端 VLM 三位一体意图识别
    # 与内嵌图片逐张 OCR 不同，slide 把整页幻灯片送 VLM，结果按页注入。
    # 识别对象是整页截图而非内嵌图，故不提取/不保留内嵌图（2026-08-09 起）：
    # 删除 MarkItDown 的内嵌图引用，返回 images=0、image_map=[]。
    slide_mode = (
        isinstance(ocr_service, SlideFusionService)
        and filepath.suffix.lower() == '.pptx'
    )
    if slide_mode:
        markdown = _strip_image_refs(markdown)
        slide_results = ocr_service.process_pptx(filepath, output_dir, verbose=verbose)
        markdown = _inject_slide_blockquotes(markdown, slide_results)
        return markdown, 0, []

    # 提取图片（PPTX/DOCX）
    if output_dir is not None:
        ext = filepath.suffix.lower()
        if ext == '.pptx':
            images_extracted, image_paths = _extract_pptx_images(filepath, output_dir)
        elif ext == '.docx':
            images_extracted, image_paths = _extract_docx_images(filepath, output_dir)

    # 按 MarkItDown 实际输出的引用名，按位置匹配构建映射
    # 映射值使用相对于 .md 文件所在目录的路径（B 内自洽）
    image_map = _build_image_map(markdown, image_paths, filepath.name)

    # 通用：删除无意义图片（md 引用 + image_map 条目 + 物理文件），所有模式统一
    markdown, image_map, image_paths = _drop_meaningless_images(
        markdown, image_map, image_paths,
    )

    # 嵌入图片识别（PPTX/DOCX）：保留原图引用，追加 blockquote 描述
    # 注意：扫描型 PDF 整页识别在前面已处理，此处仅处理嵌入图片
    if ocr_service and image_paths and filepath.suffix.lower() in {'.pptx', '.docx'}:
        if verbose:
            print(f"  开始识别 {len(image_paths)} 张图片（OCR）...")

        batch_results = ocr_service.recognize_batch(image_paths, verbose=verbose)

        # 收集有效识别结果 {提取图片名: 描述}，按位置注入到对应引用下方
        recognized = {
            img_path.name: description
            for img_path, description in batch_results.items()
            if description and not description.startswith("[")
        }
        markdown = _inject_ocr_blockquotes(markdown, image_map, recognized)

    # 注意：不在这里替换图片引用。调用方（_run_convert）负责替换，
    # 以便根据使用场景（独立 convert / pipeline）选用不同的路径策略。
    return markdown, images_extracted, image_map


def _extract_pptx_images(filepath: Path, output_dir: Path) -> tuple[int, list[Path]]:
    """从 PPTX 文件中提取所有图片，每个文件独立的图片目录"""
    try:
        from pptx import Presentation
    except ImportError:
        return 0, []

    images_dir = output_dir / f"{filepath.name}_images"
    images_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(filepath))
    image_count = 0
    image_paths = []

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext or 'png'
                    new_name = f"slide{i+1}_img{j+1}.{ext}"
                    image_path = images_dir / new_name
                    with open(image_path, 'wb') as f:
                        f.write(image_bytes)
                    image_paths.append(image_path)
                    image_count += 1
                except Exception as e:
                    logger.debug("PPTX 图片提取失败 slide%d shape%d: %s", i + 1, j + 1, e)

    return image_count, image_paths


def _extract_docx_images(filepath: Path, output_dir: Path) -> tuple[int, list[Path]]:
    """从 DOCX 文件中提取所有图片，每个文件独立的图片目录

    返回顺序按文件名中的数字部分排序（image1, image2, ..., image10, image11），
    而不是字典序（会得到 image1, image10, image11, image2, ...）。
    Word 写入 media 时通常用 imageN.png 自然命名，数字序与文档流顺序基本一致。
    """
    import re
    import zipfile

    images_dir = output_dir / f"{filepath.name}_images"
    images_dir.mkdir(parents=True, exist_ok=True)

    image_count = 0
    image_paths = []
    try:
        with zipfile.ZipFile(filepath) as zf:
            media_files = [f for f in zf.namelist() if f.startswith('word/media/')]
            media_files.sort(key=_image_sort_key)
            for media_file in media_files:
                image_data = zf.read(media_file)
                image_name = Path(media_file).name
                image_path = images_dir / image_name
                with open(image_path, 'wb') as f:
                    f.write(image_data)
                image_paths.append(image_path)
                image_count += 1
    except Exception as e:
        logger.warning("DOCX 图片提取失败: %s (%s)", filepath.name, e)

    return image_count, image_paths


def _image_sort_key(name: str):
    """图片文件名排序键：优先按数字部分，回退到字典序

    examples:
      "word/media/image1.png"  → (0, 1, "image1.png")
      "word/media/image10.png" → (0, 10, "image10.png")
      "word/media/cover.jpg"   → (1, 0, "cover.jpg")  # 无数字，字典序兜底
    """
    import re
    basename = Path(name).name
    m = re.search(r'(\d+)', basename)
    if m:
        return (0, int(m.group(1)), basename)
    return (1, 0, basename)


def _build_image_map(markdown: str, image_paths: list[Path], source_name: str) -> list[tuple[str, str]]:
    """扫描 Markdown 中实际图片引用，按位置与提取的图片路径匹配

    返回 [(ref_name, new_path), ...] 列表而非 dict：MarkItDown 对全中文
    shape.name 执行 re.sub(r"\\W", "", name) 后引用名会退化为 ".jpg"（中文字符
    全部被过滤），导致大量图片引用名重复。dict 会把重复引用合并成一条，仅剩
    一张图片能映射；list 保留重复条目，第 i 条引用 → 第 i 张提取图片（位置匹配）。

    限制：依赖"MarkItDown 输出的图片引用顺序 ≡ 我们从文档提取的图片顺序"。
    PPTX 用 slide→shape 遍历（与 MarkItDown 一致）；
    DOCX 用 image{N}.png 数字序（见 _image_sort_key）。
    若文档使用非数字命名或多图布局复杂，映射可能错位——
    彻底修复需用 python-docx 解析 <a:blip r:embed> 文档流顺序。
    """
    import re

    if not image_paths:
        return []

    img_refs = re.findall(r'!\[.*?\]\(([^)]+)\)', markdown)
    local_refs = [r for r in img_refs if not r.startswith(('http://', 'https://'))]

    image_map = []
    for i, ref in enumerate(local_refs):
        if i < len(image_paths):
            image_map.append((ref, f"{source_name}_images/{image_paths[i].name}"))

    return image_map


def _drop_meaningless_images(
    markdown: str,
    image_map: list[tuple[str, str]],
    image_paths: list[Path],
) -> tuple[str, list[tuple[str, str]], list[Path]]:
    """删除无意义图片：md 引用 + image_map 条目 + 物理文件

    复用 ImageFilter.should_recognize 判定（纯色 / 过小 / 低分辨率）。
    位置匹配：image_paths[i] ↔ image_map[i] ↔ md 第 i 个本地引用。
    对命中图片：删除 md 引用、从 image_map 移除、unlink 物理文件。
    """
    try:
        from doc_knowledge.vision import ImageFilter
    except ImportError:
        return markdown, image_map, image_paths

    filter_ = ImageFilter()
    drop_indices = {
        i for i, p in enumerate(image_paths) if not filter_.should_recognize(p)[0]
    }
    if not drop_indices:
        return markdown, image_map, image_paths

    # 1) 删除 md 中命中的本地图片引用（按位置计数，外部链接不参与）
    import re

    ref_re = re.compile(r'!\[.*?\]\(([^)]+)\)')
    local_index = 0

    def _drop(match):
        nonlocal local_index
        ref = match.group(1)
        if ref.startswith(('http://', 'https://')):
            return match.group(0)  # 外部链接不参与位置匹配
        idx = local_index
        local_index += 1
        return "" if idx in drop_indices else match.group(0)

    markdown = ref_re.sub(_drop, markdown)

    # 2) 重建 image_map，跳过被删图片
    new_image_map = [e for i, e in enumerate(image_map) if i not in drop_indices]

    # 3) 删除物理文件 + 重建 image_paths
    for i in drop_indices:
        image_paths[i].unlink(missing_ok=True)
    new_image_paths = [p for i, p in enumerate(image_paths) if i not in drop_indices]

    return markdown, new_image_map, new_image_paths


def _strip_image_refs(markdown: str) -> str:
    """删除 Markdown 中所有图片引用（![](...)）

    slide 模式识别对象是整页截图而非内嵌图，内嵌图引用是噪音，
    整页理解块已含视觉信息，无需保留内嵌图占位。
    """
    import re

    return re.sub(r'!\[.*?\]\([^)]+\)', '', markdown)


def _filter_meaningless_images(image_paths: list[Path]) -> list[Path]:
    """过滤低价值图片（文件过小 / 分辨率过低 / 纯色）

    复用 vision.py 的 ImageFilter，避免重复实现过滤逻辑。
    在 OCR 批处理前调用，减少无效的云端 API / 本地识别调用。
    """
    try:
        from doc_knowledge.vision import ImageFilter
    except ImportError:
        return image_paths

    filter_ = ImageFilter()
    return [p for p in image_paths if filter_.should_recognize(p)[0]]


def _inject_ocr_blockquotes(
    markdown: str,
    image_map: list[tuple[str, str]],
    recognized: dict[str, str],
) -> str:
    """按位置将 OCR 识别结果以 blockquote 注入到对应图片引用下方

    image_map 与 markdown 中的本地图片引用按位置一一对应（第 i 条引用 ↔ image_map[i]）。
    对已识别（recognized 中）的图片，在其引用后追加 blockquote；其余引用原样保留，
    路径替换仍由调用方（_run_convert）统一完成。

    recognized: {提取图片名: 识别描述}

    用 re.sub 一次性从左到右处理全部引用：插入的 blockquote 文本不会被重新扫描，
    避免"替换文本自身含 .jpg 锚点导致反复匹配同一位置、识别结果全部堆叠在
    第一张图片处"的问题。
    """
    import re

    if not recognized:
        return markdown

    ref_re = re.compile(r'!\[.*?\]\(([^)]+)\)')
    local_index = 0

    def _inject(match):
        nonlocal local_index
        ref = match.group(1)
        if ref.startswith(('http://', 'https://')):
            return match.group(0)  # 外部链接不参与位置匹配
        idx = local_index
        local_index += 1
        if idx < len(image_map):
            _, new_path = image_map[idx]
            desc = recognized.get(Path(new_path).name)
            if desc:
                return f"{match.group(0)}\n\n> 📷 **图片识别**: {desc}"
        return match.group(0)

    return ref_re.sub(_inject, markdown)


def _slide_desc(desc: str) -> str:
    """识别失败的错误占位符降级为可读提示

    错误占位符形态（[图片识别失败: ...] / [图片解析失败: ...]）不直接写进
    markdown，避免暴露 HTTP 错误堆栈。缺失 body 的标记 Markdown 仍可读，
    原样返回（不丢弃）。
    """
    if desc.startswith(("[图片识别失败", "[图片解析失败")):
        return "⚠️ 本页图表语义识别失败（VLM 请求异常），原始文字已保留"
    return desc


def _parse_slide_result(text: str) -> Optional[dict]:
    """解析 VLM 输出的约定标记 Markdown（[DK-xxx] 标签）；缺/空 body 返回 None

    VLM 文本形态：
    - 含 [DK-正文] 且非空 → dict（title/overview/structure 可选，body 必需）
    - 缺 [DK-正文] 或正文为空 → None（降级，保留原文）
    - 失败占位符（[图片识别失败...] / [图片解析失败...]）→ None

    标签值 = 标签后到下一个 [DK- 标签（或文本末尾）前的全部内容；
    [DK-正文] 特殊——取到文本末尾，故 body 可含任意 Markdown 不被后续标签截断。
    纯文本标记无转义负担，无需 json.loads 容错层（2026-08-10 JSON→标记迁移）。
    """
    import re

    if not text or not text.strip() or text.startswith(("[图片识别失败", "[图片解析失败")):
        return None
    result: dict = {}
    # 可选字段：标题/概述/结构 —— 标签后到下一标签前为值
    for key, tag in (("title", "标题"), ("overview", "概述"), ("structure", "结构")):
        m = re.search(rf"\[DK-{tag}\](.*?)(?=\[DK-|\Z)", text, re.S)
        if m:
            result[key] = m.group(1).strip()
    # 必需字段：正文 —— 标签后到文本末尾
    m = re.search(r"\[DK-正文\](.*)\Z", text, re.S)
    if not m or not m.group(1).strip():
        return None
    result["body"] = m.group(1).strip()
    return result


def _render_slide_page(page_text: str, result: dict) -> str:
    """渲染单页结构化输出：📊 标题（结构）+ 概述 + body + markitdown 原文折叠

    result: _parse_slide_result 解析后的 dict（title/overview/structure/body）
    page_text: MarkItDown 该页原始文本（<details> 折叠作兜底参考，可靠性 > 体积）
    """
    import json

    def _fmt(v):
        """字段非 str（VLM 嵌套 dict/list）时转为可读 Markdown，避免裸 JSON

        - dict：键 → `###` 小标题，值（str 或 list）→ 项目符号列表
        - list：→ 项目符号列表
        - 其他标量 → JSON 序列化兜底
        """
        if isinstance(v, str):
            return v
        if isinstance(v, dict):
            blocks = []
            for k, val in v.items():
                items = val if isinstance(val, list) else [val]
                blocks.append(
                    f"### {k}\n" + "\n".join(f"- {_fmt(x)}" for x in items)
                )
            return "\n\n".join(blocks)
        if isinstance(v, list):
            return "\n".join(f"- {_fmt(x)}" for x in v)
        return json.dumps(v, ensure_ascii=False) if v is not None else ""

    title = _fmt(result.get("title")).strip()
    structure = _fmt(result.get("structure")).strip()
    overview = _fmt(result.get("overview")).strip()
    body = _fmt(result.get("body")).strip()

    lines = []
    if title:
        head = f"📊 **{title}**"
        if structure:
            head += f"（{structure}）"
        lines.append(head)
    if overview:
        lines.append(f"**概述**：{overview}")
    if body:
        lines.append(body)
    page_text = page_text.strip()
    if page_text:
        lines.append(
            f"<details><summary>markitdown 原文（兜底参考）</summary>\n{page_text}\n</details>"
        )
    return "\n\n".join(lines)


def _inject_slide_blockquotes(markdown: str, slide_results: dict[int, str]) -> str:
    """按 <!-- Slide number: N --> 边界，每页注入结构化输出

    slide_results: {页码: VLM 文本}（页码与 MarkItDown 输出的 Slide number 对应）。
    VLM 文本分派：
    - 合法结构化 JSON → _render_slide_page 渲染（标题（结构）+ 概述 + body + 原文折叠）
    - 非 JSON 文本 → 降级 `> 📊 **整页理解**: {文本}`（保留原文，兼容旧产物）
    - 失败占位符（[图片识别失败...]）→ 降级提示
    无结果的页（页码缺失或识别失败）保持原样。
    """
    import re

    if not slide_results:
        return markdown

    # 按页标记切分（lookahead 保留标记本身），逐页渲染后重组
    parts = re.split(r'(?=<!-- Slide number: \d+ -->)', markdown)
    out = []
    for part in parts:
        m = re.match(r'<!-- Slide number: (\d+) -->', part)
        if not m:
            out.append(part)
            continue
        num = int(m.group(1))
        desc = slide_results.get(num)
        if not desc:
            out.append(part)
            continue
        result = _parse_slide_result(desc)
        if result is not None:
            page_text = part[m.end():]
            out.append(f"<!-- Slide number: {num} -->\n\n{_render_slide_page(page_text, result)}")
        else:
            out.append(f"{part.rstrip()}\n\n> 📊 **整页理解**: {_slide_desc(desc)}")
    # markdown 以 <!-- Slide number: 1 --> 开头时 parts[0] 为空串，join 会拼出前导 \n
    return "\n".join(out).lstrip("\n")


def _update_slide_blockquotes(markdown: str, slide_results: dict[int, str]) -> str:
    """补跑后更新已有 markdown 的指定页结构化输出（retry-slide 用）

    - 目标页（失败页）：剥离旧失败块（新形态独立 ⚠️ 行 / 旧形态 blockquote），
      用新结果重新渲染——成功→结构化输出（含原文折叠），仍失败→降级提示
    - 非目标页：保持原样（不重复注入、不覆盖已有成功块）
    """
    import re

    if not slide_results:
        return markdown

    target = set(slide_results)
    fail_re = re.compile(
        r"⚠️ 本页图表语义识别失败（VLM 请求异常），原始文字已保留"
        r"|> 📊 \*\*整页理解\*\*: (?:⚠️|\[[^\r\n]*\])"
    )

    parts = re.split(r'(?=<!-- Slide number: \d+ -->)', markdown)
    out = []
    for part in parts:
        m = re.match(r'<!-- Slide number: (\d+) -->', part)
        if not m or int(m.group(1)) not in target:
            out.append(part)
            continue
        num = int(m.group(1))
        page_text = fail_re.sub("", part[m.end():])
        desc = slide_results.get(num)
        result = _parse_slide_result(desc) if desc else None
        if result is not None:
            out.append(
                f"<!-- Slide number: {num} -->\n\n{_render_slide_page(page_text, result)}"
            )
        else:
            # _slide_desc 对 [ 开头占位符已返回带 ⚠️ 的提示，此处不再加前缀，避免双重 ⚠️
            out.append(
                f"<!-- Slide number: {num} -->\n\n{_slide_desc(desc)}\n\n{page_text.strip()}"
            )
    return "\n".join(out).lstrip("\n")


def get_supported_extensions() -> list[str]:
    """获取支持的文件扩展名列表"""
    return [
        ".pdf", ".docx", ".pptx", ".xlsx", ".xls",
        ".html", ".htm", ".epub", ".csv",
        ".txt", ".md",
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp",
        ".mp3", ".wav", ".m4a",
        ".zip", ".msg", ".ipynb",
    ]
